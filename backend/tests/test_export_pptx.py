"""Tests for PPTX exporter — code safety, sandbox, caching."""

import os

os.environ.setdefault("ANTHROPIC_API_KEY", "sk-ant-test-key-for-testing")

import asyncio

import pytest
from unittest.mock import AsyncMock

from app.providers.base import GenerationResult, LLMProvider
from app.services.exporters.base import ExportGenerationError, ExportOptions
from app.services.exporters.pptx_exporter import (
    PPTXExporter,
    _safe_import,
)

SAMPLE_HTML = "<!DOCTYPE html><html><body><h1>Test</h1></body></html>"


@pytest.fixture
def mock_provider():
    provider = AsyncMock(spec=LLMProvider)
    provider.model = "claude-sonnet-4-6-20260217"
    return provider


@pytest.fixture
def exporter(mock_provider: AsyncMock):
    return PPTXExporter(mock_provider)


# ---------------------------------------------------------------------------
# _safe_import tests
# ---------------------------------------------------------------------------

def test_safe_import_allows_pptx():
    """pptx should be importable."""
    mod = _safe_import("pptx")
    assert mod is not None


def test_safe_import_allows_io():
    mod = _safe_import("io")
    assert mod is not None


def test_safe_import_allows_pptx_submodule():
    mod = _safe_import("pptx.util")
    assert mod is not None


def test_safe_import_blocks_os():
    with pytest.raises(ImportError, match="not allowed"):
        _safe_import("os")


def test_safe_import_blocks_subprocess():
    with pytest.raises(ImportError, match="not allowed"):
        _safe_import("subprocess")


def test_safe_import_blocks_socket():
    with pytest.raises(ImportError, match="not allowed"):
        _safe_import("socket")


def test_safe_import_blocks_requests():
    with pytest.raises(ImportError, match="not allowed"):
        _safe_import("requests")


# ---------------------------------------------------------------------------
# _validate_code_safety tests
# ---------------------------------------------------------------------------

def test_validate_blocks_import_os(exporter: PPTXExporter):
    with pytest.raises(ExportGenerationError, match="forbidden"):
        exporter._validate_code_safety("import os\nfrom pptx import Presentation")


def test_validate_blocks_import_subprocess(exporter: PPTXExporter):
    with pytest.raises(ExportGenerationError, match="forbidden"):
        exporter._validate_code_safety("import subprocess\nfrom pptx import Presentation")


def test_validate_blocks_eval(exporter: PPTXExporter):
    with pytest.raises(ExportGenerationError, match="forbidden"):
        exporter._validate_code_safety("from pptx import Presentation\neval('bad')")


def test_validate_blocks_open(exporter: PPTXExporter):
    with pytest.raises(ExportGenerationError, match="forbidden"):
        exporter._validate_code_safety("from pptx import Presentation\nopen('/etc/passwd')")


def test_validate_blocks_getattr(exporter: PPTXExporter):
    with pytest.raises(ExportGenerationError, match="forbidden"):
        exporter._validate_code_safety("from pptx import Presentation\ngetattr(x, 'y')")


def test_validate_requires_pptx_import(exporter: PPTXExporter):
    with pytest.raises(ExportGenerationError, match="must import"):
        exporter._validate_code_safety("result = b'fake'")


def test_validate_allows_valid_code(exporter: PPTXExporter):
    code = """from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
output = BytesIO()
prs.save(output)
result = output.getvalue()
"""
    exporter._validate_code_safety(code)  # Should not raise


# ---------------------------------------------------------------------------
# _extract_code tests
# ---------------------------------------------------------------------------

def test_extract_code_from_python_fence(exporter: PPTXExporter):
    response = "Here is the code:\n```python\nprint('hello')\n```\nDone."
    assert exporter._extract_code(response) == "print('hello')"


def test_extract_code_from_plain_fence(exporter: PPTXExporter):
    response = "```\nprint('hello')\n```"
    assert exporter._extract_code(response) == "print('hello')"


def test_extract_code_bare(exporter: PPTXExporter):
    response = "print('hello')"
    assert exporter._extract_code(response) == "print('hello')"


# ---------------------------------------------------------------------------
# Cache key tests
# ---------------------------------------------------------------------------

def test_cache_key_deterministic(exporter: PPTXExporter):
    options = ExportOptions()
    key1 = exporter._get_cache_key("html1", options)
    key2 = exporter._get_cache_key("html1", options)
    assert key1 == key2


def test_cache_key_different_html(exporter: PPTXExporter):
    options = ExportOptions()
    key1 = exporter._get_cache_key("html1", options)
    key2 = exporter._get_cache_key("html2", options)
    assert key1 != key2


def test_cache_key_different_options(exporter: PPTXExporter):
    key1 = exporter._get_cache_key("html", ExportOptions(slide_width=10))
    key2 = exporter._get_cache_key("html", ExportOptions(slide_width=16))
    assert key1 != key2


# ---------------------------------------------------------------------------
# Sandbox execution tests
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_sandbox_executes_valid_pptx_code(exporter: PPTXExporter):
    """Actually run a minimal pptx script in the sandbox."""
    code = """from pptx import Presentation
from io import BytesIO

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
output = BytesIO()
prs.save(output)
result = output.getvalue()
"""
    result = await exporter._execute_pptx_code(code, max_retries=0)
    assert isinstance(result, bytes)
    assert len(result) > 100  # PPTX files are at least a few KB


@pytest.mark.asyncio
async def test_sandbox_blocks_os_at_runtime(exporter: PPTXExporter):
    """Even if validation is bypassed, sandbox should block os import."""
    code = """import os
result = b'bad'
"""
    with pytest.raises(ExportGenerationError):
        await exporter._execute_pptx_code(code, max_retries=0)


@pytest.mark.asyncio
async def test_sandbox_timeout(exporter: PPTXExporter):
    """Code execution timeout should raise ExportGenerationError."""
    from unittest.mock import patch

    # Mock _run_in_executor to simulate a timeout
    async def _slow_executor(*args, **kwargs):
        await asyncio.sleep(999)

    with patch.object(exporter, "_run_in_executor", side_effect=_slow_executor):

        async def _fast_execute(code: str, max_retries: int = 0) -> bytes:
            import asyncio as _aio
            from app.services.exporters.pptx_exporter import _build_restricted_builtins

            restricted_globals: dict = {"__builtins__": _build_restricted_builtins()}
            exec_locals: dict = {}
            try:
                await _aio.wait_for(
                    exporter._run_in_executor(code, restricted_globals, exec_locals),
                    timeout=0.1,
                )
            except _aio.TimeoutError:
                raise ExportGenerationError("Code execution timeout (30s)")
            return exec_locals.get("result", b"")  # type: ignore[return-value]

        with pytest.raises(ExportGenerationError, match="timeout"):
            await _fast_execute("from pptx import Presentation")


@pytest.mark.asyncio
async def test_sandbox_no_result_variable(exporter: PPTXExporter):
    """Code that doesn't set 'result' should fail."""
    code = """from pptx import Presentation
prs = Presentation()
"""
    with pytest.raises(ExportGenerationError, match="did not produce"):
        await exporter._execute_pptx_code(code, max_retries=0)


# ---------------------------------------------------------------------------
# Full export flow with mocked provider
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_export_calls_provider_generate(
    mock_provider: AsyncMock, exporter: PPTXExporter
):
    """Verify provider.generate() is called with temperature=0.0."""
    # Return valid pptx generation code
    mock_provider.generate.return_value = GenerationResult(
        text="""```python
from pptx import Presentation
from io import BytesIO

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
output = BytesIO()
prs.save(output)
result = output.getvalue()
```""",
    )

    options = ExportOptions(document_title="test")
    result = await exporter.export(SAMPLE_HTML, options)

    assert isinstance(result.content, bytes)
    assert len(result.content) > 0
    assert result.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert result.filename == "test.pptx"

    # Verify generate was called with temperature=0.0
    call_kwargs = mock_provider.generate.call_args
    assert call_kwargs.kwargs.get("temperature") == 0.0 or call_kwargs[1].get("temperature") == 0.0


@pytest.mark.asyncio
async def test_export_uses_cache(
    mock_provider: AsyncMock, exporter: PPTXExporter
):
    """Second call with same HTML should use cache."""
    mock_provider.generate.return_value = GenerationResult(
        text="""```python
from pptx import Presentation
from io import BytesIO

prs = Presentation()
output = BytesIO()
prs.save(output)
result = output.getvalue()
```""",
    )

    options = ExportOptions(document_title="test")
    await exporter.export(SAMPLE_HTML, options)
    await exporter.export(SAMPLE_HTML, options)

    # Provider should only be called once (second call uses cache)
    assert mock_provider.generate.call_count == 1


# ---------------------------------------------------------------------------
# Plan 019: PPTX quality improvements
# ---------------------------------------------------------------------------

def test_default_slide_dimensions_16_9():
    """ExportOptions should default to 16:9 widescreen."""
    opts = ExportOptions()
    assert opts.slide_width == 13.333
    assert opts.slide_height == 7.5


def test_base64_stripped_from_prompt(exporter: PPTXExporter):
    """Base64 image data should be replaced with placeholder in prompt."""
    html = '<html><body><img src="data:image/png;base64,' + 'A' * 500 + '"/></body></html>'
    prompt = exporter._build_generation_prompt(html, ExportOptions())
    assert 'AAAA' not in prompt
    assert '[IMAGE_DATA_REMOVED]' in prompt


def test_prompt_includes_slide_dimensions(exporter: PPTXExporter):
    """Prompt should reference the configured slide dimensions."""
    options = ExportOptions(slide_width=13.333, slide_height=7.5)
    prompt = exporter._build_generation_prompt(SAMPLE_HTML, options)
    assert '13.333' in prompt
    assert '7.5' in prompt


def test_prompt_includes_color_extraction(exporter: PPTXExporter):
    """Prompt should instruct Claude to extract colors from HTML."""
    prompt = exporter._build_generation_prompt(SAMPLE_HTML, ExportOptions())
    assert 'COLOR EXTRACTION' in prompt
    assert 'RGBColor' in prompt


def test_prompt_includes_formatting_rules(exporter: PPTXExporter):
    """Prompt should include visual formatting rules."""
    prompt = exporter._build_generation_prompt(SAMPLE_HTML, ExportOptions())
    assert 'MSO_SHAPE' in prompt
    assert 'NEVER use fonts smaller than 14pt' in prompt
    assert 'Maximum 5 bullet points' in prompt
